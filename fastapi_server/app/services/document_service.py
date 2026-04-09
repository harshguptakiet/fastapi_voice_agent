import io
from pathlib import Path

from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from app.services.object_storage_service import object_storage

class DocumentService:
    def __init__(self, storage_dir: str = "storage/documents"):
        self.storage_dir = Path(storage_dir)
        self.storage_dir.mkdir(parents=True, exist_ok=True)

    async def save_and_extract(self, file_name: str, content: bytes) -> dict:
        # 1. Save raw file
        raw_uri = object_storage.save_raw(file_name, content)

        # 2. Extract text
        ext = Path(file_name).suffix.lower()
        extracted_text = ""

        try:
            if ext == ".pdf":
                extracted_text = self._extract_pdf(content)
            elif ext == ".docx":
                extracted_text = self._extract_docx(content)
            elif ext == ".pptx":
                extracted_text = self._extract_pptx(content)
            else:
                raise ValueError(f"Unsupported file extension: {ext}")
        except Exception as e:
            # Propagate extraction errors so callers can handle them explicitly,
            # instead of treating error messages as extracted content.
            raise RuntimeError(f"Error extracting text: {e}") from e

        # 3. Clean and save extracted text
        cleaned_text = self._clean_text(extracted_text)
        text_file_name = f"{Path(file_name).stem}.txt"
        text_uri = object_storage.save_text(text_file_name, cleaned_text)

        return {
            "file_name": file_name,
            "raw_path": raw_uri,
            "text_path": text_uri,
            "content_preview": cleaned_text[:200] + "..." if len(cleaned_text) > 200 else cleaned_text,
            "extracted_text": cleaned_text,
        }

    def list_documents(self) -> list[str]:
        return object_storage.list_text_documents()

    def _extract_pdf(self, content: bytes) -> str:
        reader = PdfReader(io.BytesIO(content))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

    def _extract_docx(self, content: bytes) -> str:
        doc = Document(io.BytesIO(content))
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def _extract_pptx(self, content: bytes) -> str:
        prs = Presentation(io.BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _clean_text(self, text: str) -> str:
        # Basic cleaning: remove excessive whitespace
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)

document_service = DocumentService()
